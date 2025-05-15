from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

class PPTGenerator:
    def __init__(self):
        self.prs = Presentation()
        # Set default slide size to 16:9 (widescreen)
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        
        # Color schemes
        self.color_schemes = {
            "professional": {
                "primary": RGBColor(0, 85, 140),    # Deep Blue
                "secondary": RGBColor(225, 225, 225),  # Light Grey
                "accent": RGBColor(0, 164, 239),    # Bright Blue
                "text": RGBColor(51, 51, 51)        # Dark Grey
            },
            "modern": {
                "primary": RGBColor(52, 73, 94),    # Dark Slate
                "secondary": RGBColor(236, 240, 241),  # Cloud
                "accent": RGBColor(26, 188, 156),   # Turquoise
                "text": RGBColor(44, 62, 80)        # Midnight Blue
            },
            "creative": {
                "primary": RGBColor(155, 89, 182),  # Purple
                "secondary": RGBColor(250, 250, 250),  # White Smoke
                "accent": RGBColor(241, 196, 15),   # Yellow
                "text": RGBColor(52, 73, 94)        # Dark Slate
            }
        }
        
        # Default color scheme
        self.current_scheme = self.color_schemes["professional"]
    
    def set_color_scheme(self, scheme_name):
        """Set the color scheme for the presentation"""
        if scheme_name in self.color_schemes:
            self.current_scheme = self.color_schemes[scheme_name]
            return True
        return False
    
    def create_title_slide(self, name, title, tagline=None):
        """Create the title slide with name and professional title"""
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.current_scheme["secondary"]
        
        # Add title (name)
        title_shape = slide.shapes.title
        title_shape.text = name
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_run = title_para.runs[0]
        title_run.font.size = Pt(54)
        title_run.font.color.rgb = self.current_scheme["primary"]
        title_run.font.bold = True
        
        # Add subtitle (professional title)
        subtitle = slide.placeholders[1]
        subtitle.text = title
        subtitle_para = subtitle.text_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_run = subtitle_para.runs[0]
        subtitle_run.font.size = Pt(32)
        subtitle_run.font.color.rgb = self.current_scheme["text"]
        
        # Add tagline if provided
        if tagline:
            subtitle.text += f"\n\n{tagline}"
            tagline_para = subtitle.text_frame.paragraphs[2]
            tagline_para.alignment = PP_ALIGN.CENTER
            tagline_run = tagline_para.runs[0]
            tagline_run.font.size = Pt(24)
            tagline_run.font.color.rgb = self.current_scheme["accent"]
            tagline_run.font.italic = True
        
        return slide
    
    def create_content_slide(self, title, bullet_points=None, layout_type="bullet"):
        """Create a content slide with title and bullet points"""
        slide_layout = self.prs.slide_layouts[1]  # Content slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.current_scheme["secondary"]
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_run = title_para.runs[0]
        title_run.font.size = Pt(40)
        title_run.font.color.rgb = self.current_scheme["primary"]
        title_run.font.bold = True
        
        # Add content based on layout type
        if layout_type == "bullet" and bullet_points:
            content = slide.placeholders[1]
            text_frame = content.text_frame
            text_frame.clear()  # Clear existing text
            
            for point in bullet_points:
                p = text_frame.add_paragraph()
                p.text = point
                p.level = 0
                p.alignment = PP_ALIGN.LEFT
                run = p.runs[0]
                run.font.size = Pt(24)
                run.font.color.rgb = self.current_scheme["text"]
        
        return slide
    
    def create_experience_slide(self, experiences):
        """Create a slide with work experience details"""
        for exp in experiences:
            title = f"Experience: {exp.get('title')} at {exp.get('company')}"
            bullets = [
                f"{exp.get('dates', 'N/A')}",
                *exp.get('responsibilities', [])
            ]
            self.create_content_slide(title, bullets)
    
    def create_education_slide(self, education_items):
        """Create a slide with education details"""
        title = "Education"
        bullets = []
        
        for edu in education_items:
            bullet = f"{edu.get('degree')} - {edu.get('institution')}, {edu.get('year')}"
            if edu.get('gpa'):
                bullet += f" (GPA: {edu.get('gpa')})"
            bullets.append(bullet)
            
            # Add any achievements as sub-bullets
            if edu.get('achievements'):
                for achievement in edu.get('achievements'):
                    bullets.append(f"  â€¢ {achievement}")
        
        self.create_content_slide(title, bullets)
    
    def create_skills_slide(self, skills):
        """Create a slide with categorized skills"""
        title = "Skills & Expertise"
        
        # Create a custom slide to display categorized skills
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.current_scheme["secondary"]
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_run = title_para.runs[0]
        title_run.font.size = Pt(40)
        title_run.font.color.rgb = self.current_scheme["primary"]
        title_run.font.bold = True
        
        # Add content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Add each skill category
        categories = ["technical", "soft", "domain"]
        for category in categories:
            if category in skills and skills[category]:
                p = text_frame.add_paragraph()
                p.text = f"{category.capitalize()} Skills:"
                p.alignment = PP_ALIGN.LEFT
                run = p.runs[0]
                run.font.size = Pt(24)
                run.font.color.rgb = self.current_scheme["accent"]
                run.font.bold = True
                
                # Add skills in this category
                for skill in skills[category]:
                    p = text_frame.add_paragraph()
                    p.text = f"â€¢ {skill}"
                    p.level = 1
                    p.alignment = PP_ALIGN.LEFT
                    run = p.runs[0]
                    run.font.size = Pt(20)
                    run.font.color.rgb = self.current_scheme["text"]
                
                # Add empty paragraph as spacer
                text_frame.add_paragraph()
        
        return slide
    
    def create_contact_slide(self, contact_info):
        """Create a contact information slide"""
        title = "Contact Information"
        bullets = [
            f"Email: {contact_info.get('email', 'N/A')}",
            f"Phone: {contact_info.get('phone', 'N/A')}",
            f"LinkedIn: {contact_info.get('linkedin', 'N/A')}",
            f"Portfolio: {contact_info.get('portfolio', 'N/A')}"
        ]
        
        self.create_content_slide(title, bullets)
    
    def generate_from_ai_content(self, ai_content, color_scheme="professional"):
        """Generate a complete PowerPoint from AI-generated content"""
        # Set color scheme
        self.set_color_scheme(color_scheme)
        
        # Create title slide
        title_data = ai_content.get("title_slide", {})
        self.create_title_slide(
            title_data.get("name", ""),
            title_data.get("title", ""),
            title_data.get("tagline", "")
        )
        
        # Create About Me slide
        about_data = ai_content.get("about_me", {})
        self.create_content_slide(
            "About Me",
            about_data.get("points", [])
        )
        
        # Create Experience slides
        experiences = ai_content.get("work_experience", [])
        self.create_experience_slide(experiences)
        
        # Create Education slide
        education = ai_content.get("education", [])
        self.create_education_slide(education)
        
        # Create Skills slide
        skills = ai_content.get("skills", {})
        self.create_skills_slide(skills)
        
        # Create Achievements slide
        achievements = ai_content.get("achievements", [])
        if achievements:
            self.create_content_slide("Achievements & Certifications", achievements)
        
        # Create Contact slide
        contact = ai_content.get("contact", {})
        self.create_contact_slide(contact)
    
    def save(self, filename="resume_presentation.pptx"):
        """Save the presentation to file"""
        self.prs.save(filename)
        return filename